from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from models.model import PresentationContent
import io
import os
import base64
import tempfile


def generate_presentation_pptx(content: PresentationContent, output_buffer: io.BytesIO):
    """
    Generates a PPTX file based on the structured content provided by the agent.
    This version includes dynamic image placement and modern bullet point styling.
    """
    prs = Presentation() # Starts with a default, blank presentation

    # Set presentation dimensions to 16:9 widescreen for modern look
    prs.slide_width = Inches(13.333) # Equivalent to 16:9 aspect ratio
    prs.slide_height = Inches(7.5) # Equivalent to 16:9 aspect ratio

    # --- Define a simple "theme" based on overall_theme ---
    # In a real application, this would be a more sophisticated mapping to templates or styles.
    # For now, we'll just define some default colors and potentially adjust based on a theme hint.
    text_color_primary = RGBColor(0, 0, 0) # Black
    text_color_secondary = RGBColor(50, 50, 50) # Dark gray
    accent_color = RGBColor(0, 112, 192) # Default Blue

    if content.overall_theme:
        theme_lower = content.overall_theme.lower()
        print(f"Applying theme: {theme_lower}")
        if "professional" in theme_lower or "corporate" in theme_lower:
            text_color_primary = RGBColor(30, 30, 30)
            text_color_secondary = RGBColor(80, 80, 80)
            accent_color = RGBColor(68, 114, 196) # Darker blue
        elif "energetic" in theme_lower or "dynamic" in theme_lower:
            accent_color = RGBColor(255, 120, 0) # Orange
        elif "minimalist" in theme_lower:
            text_color_primary = RGBColor(10, 10, 10)
            text_color_secondary = RGBColor(60, 60, 60)
            accent_color = RGBColor(150, 150, 150) # Lighter gray accent

    for i, slide_data in enumerate(content.slides):
        # Always start with a blank slide layout for manual placement
        slide_layout = prs.slide_layouts[6] # Index 6 is usually 'Blank'
        slide = prs.slides.add_slide(slide_layout)

        # Define common margins and sizes
        left_margin = Inches(0.75)
        right_margin = Inches(0.75)
        top_margin = Inches(0.5)
        bottom_margin = Inches(0.5)
        content_width = prs.slide_width - left_margin - right_margin
        content_height = prs.slide_height - top_margin - bottom_margin

        # --- Add Title ---
        title_box = slide.shapes.add_textbox(left_margin, top_margin, content_width, Inches(1))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = slide_data.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = text_color_primary
        p.alignment = PP_ALIGN.LEFT

        # --- Determine Layout based on content presence ---
        has_bullets = bool(slide_data.bullet_points)
        has_image = bool(slide_data.image_description)
        has_image_base64 = hasattr(slide_data, "image_base64") and bool(slide_data.image_base64)

        # Initial content area for text/bullets
        text_top = top_margin + Inches(1.2)
        text_height = content_height - Inches(1.2) # Remaining height after title
        current_left_offset = left_margin
        current_text_width = content_width

        # --- Dynamic Image Placement and Content Area Adjustment ---
        if has_image:
            image_placeholder_width = Inches(4) # Fixed width for placeholder for now
            image_placeholder_height = Inches(3) # Fixed height

            # Layout 1: Title + Image on Right + Bullets on Left
            if has_bullets:
                current_text_width = content_width - image_placeholder_width - Inches(0.5) # Space for image and margin
                image_left = left_margin + current_text_width + Inches(0.5)
                image_top = text_top + (text_height - image_placeholder_height) / 2 # Vertically center image in remaining space
            # Layout 2: Title + Image (Centered)
            else:
                current_text_width = 0 # No text box if only image
                image_left = (prs.slide_width - image_placeholder_width) / 2
                image_top = text_top + (text_height - image_placeholder_height) / 2 # Centered below title


            
            if has_image_base64:
                # Decode base64 and save to temp file
                

                img_bytes = base64.b64decode(slide_data.image_base64)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                    tmp_img.write(img_bytes)
                    tmp_img_path = tmp_img.name

                # Add image to slide
                slide.shapes.add_picture(
                    tmp_img_path,
                    image_left,
                    image_top,
                    width=image_placeholder_width,
                    height=image_placeholder_height
                )
                os.remove(tmp_img_path)  # Clean up temp file
            else:
            
                # Add image placeholder (enhanced for modern look)
                img_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    image_left,
                    image_top,
                    image_placeholder_width,
                    image_placeholder_height
                )
                img_shape.fill.background() # No solid fill
                img_shape.line.color.rgb = accent_color # Use accent color for border
                img_shape.line.width = Pt(2) # Thicker line
                # img_shape.line.dash_style =P #pptx.enum.shapes.MsoLineDashStyle.DASH # Dashed line
                img_shape.line.dash_style = MSO_LINE.DASH_DOT_DOT

                # Add text description over the placeholder
                text_box_img = img_shape.text_frame
                text_box_img.word_wrap = True
                text_box_img.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                text_box_img.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_img = text_box_img.paragraphs[0]
                p_img.text = f"Image Suggestion:\n'{slide_data.image_description}'"
                p_img.font.size = Pt(10)
                p_img.font.color.rgb = text_color_secondary
                p_img.alignment = PP_ALIGN.CENTER

        # --- Add Bullet Points (Modern Styling) ---
        if has_bullets:
            bullet_box = slide.shapes.add_textbox(current_left_offset, text_top, current_text_width, text_height)
            bullet_frame = bullet_box.text_frame
            bullet_frame.word_wrap = True
            bullet_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT # Adjust text size to fit shape, or vice versa
            bullet_frame.vertical_anchor = MSO_ANCHOR.TOP # Anchor text to the top of the box

            for j, point in enumerate(slide_data.bullet_points):
                p = bullet_frame.add_paragraph()
                p.text = '🔷 ' + point
                p.level = 0 # Top-level bullet

                # Modern bullet point styling: Using a square bullet character
                # and adjusting font and indentation
                p.font.size = Pt(18) # Slightly smaller font for body
                p.font.color.rgb = text_color_secondary

                # Custom bullet character and indentation
                p.space_before = Pt(5) # Small space before each bullet
                p.space_after = Pt(5)  # Small space after each bullet
                p.left_indent = Inches(0.2) # Indent for the bullet point itself
                p.first_line_indent = Inches(-0.2) # Negative indent to bring bullet symbol left of text

                # Set custom bullet symbol (square)
                # This requires setting bullet_char and bullet_font
                p.bullet_char = '▣' #'■'# Unicode square character
                # You might need to set a specific font for the bullet char if default doesn't render well
                # p.bullet_font.name = 'Arial' # Example

    # Save the presentation to the provided in-memory buffer
    prs.save(output_buffer)
    # prs.save('test.pptx')

